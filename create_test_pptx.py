from pptx import Presentation
from pptx.util import Inches, Pt

def create_test_presentation(output_path="example.pptx"):
    """Create a sample PowerPoint file with slides and titles for testing."""
    prs = Presentation()
    
    # Define slide titles
    slide_titles = [
        "Introduction to TitleWave",
        "Key Features and Benefits",
        "Technical Implementation",
        "Future Development",
        "Questions & Answers"
    ]
    
    # Create slides with titles
    for title in slide_titles:
        slide_layout = prs.slide_layouts[0]  # Title slide layout
        slide = prs.slides.add_slide(slide_layout)
        title_shape = slide.shapes.title
        title_shape.text = title
        
        # Add some placeholder content
        if slide.placeholders[1].has_text_frame:
            tf = slide.placeholders[1].text_frame
            tf.text = "Sample content for " + title
    
    # Save the presentation
    prs.save(output_path)
    print(f"Created test presentation at {output_path}")
    return output_path

if __name__ == "__main__":
    create_test_presentation() 
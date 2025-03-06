from pptx import Presentation

def extract_slide_titles(pptx_path):
    """Extract slide titles from a PowerPoint presentation."""
    prs = Presentation(pptx_path)
    
    titles = []
    for slide in prs.slides:
        # Try to get the title from the slide
        title = ""
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text:
                # For simplicity, we'll assume the first shape with text is the title
                # A more robust solution would check shape type or placeholder type
                title = shape.text
                break
        
        # If no title found, provide a placeholder
        if not title:
            title = f"Slide #{len(titles)+1} (No Title)"
        
        titles.append(title)
    
    return titles 
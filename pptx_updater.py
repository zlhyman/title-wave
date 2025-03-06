from pptx import Presentation

def update_pptx_titles(input_pptx_path, new_titles, output_pptx_path):
    """Update slide titles in a PowerPoint presentation."""
    prs = Presentation(input_pptx_path)
    
    # Ensure we have enough titles
    if len(new_titles) < len(prs.slides):
        # Pad with empty strings if needed
        new_titles.extend([""] * (len(prs.slides) - len(new_titles)))
    
    for i, slide in enumerate(prs.slides):
        if i >= len(new_titles):
            break
            
        # Find the title shape
        for shape in slide.shapes:
            # Update the first shape with text (assuming it's the title)
            if hasattr(shape, "text") and shape.text:
                shape.text = new_titles[i]
                break
    
    # Save the updated presentation
    prs.save(output_pptx_path) 